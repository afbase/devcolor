#!/usr/bin/env python3
"""Generate the OWASP Top 10:2025 workshop deck as a .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- palette ---------------------------------------------------------------
NAVY   = RGBColor(0x14, 0x1B, 0x2E)
NAVY2  = RGBColor(0x1E, 0x29, 0x44)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GREY   = RGBColor(0x5B, 0x63, 0x72)
LIGHT  = RGBColor(0xEC, 0xEF, 0xF4)
TEAL   = RGBColor(0x2E, 0xC4, 0xB6)
RED    = RGBColor(0xD1, 0x3B, 0x3B)
REDBG  = RGBColor(0xFB, 0xE9, 0xE9)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
GRNBG  = RGBColor(0xE7, 0xF2, 0xE8)
AMBER  = RGBColor(0xE8, 0x8A, 0x1A)
INK    = RGBColor(0x22, 0x28, 0x33)
CODEBG = RGBColor(0x16, 0x1B, 0x24)
CODEFG = RGBColor(0xD6, 0xE2, 0xF0)

FONT = "Arial"
MONO = "Consolas"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def slide():
    return prs.slides.add_slide(BLANK)

def bg(s, color):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = color

def box(s, x, y, w, h, fill=None, line=None, line_w=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w or 1)
    sp.shadow.inherit = False
    return sp

def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (text, size, bold, color, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after); p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, b, c, f) in para:
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.bold = b
            r.font.color.rgb = c; r.font.name = f or FONT
    return tb

def P(t, sz=18, b=False, c=INK, f=FONT):
    return [(t, sz, b, c, f)]

def header(s, code, title, accent=TEAL, kicker=None):
    """Standard content-slide header band."""
    box(s, 0, 0, 13.333, 1.15, fill=NAVY)
    box(s, 0, 1.15, 13.333, 0.06, fill=accent)
    runs = [[(code + "   ", 24, True, accent, FONT), (title, 24, True, WHITE, FONT)]]
    text(s, 0.55, 0, 12.2, 1.15, runs, anchor=MSO_ANCHOR.MIDDLE)
    if kicker:
        text(s, 0.55, 0.72, 12.2, 0.4, [[(kicker, 12, False, RGBColor(0x9F,0xB0,0xC8), FONT)]])

def codebox(s, x, y, w, h, lines, title=None):
    box(s, x, y, w, h, fill=CODEBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs = []
    if title:
        runs.append([(title, 11, True, RGBColor(0x7C,0x8A,0x9E), MONO)])
    for ln, col in lines:
        runs.append([(ln, 12.5, False, col, MONO)])
    text(s, x+0.25, y+0.18, w-0.5, h-0.36, runs, space_after=3, line_spacing=1.05)

def chip(s, x, y, w, label, fill, fg=WHITE):
    box(s, x, y, w, 0.42, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, x, y, w, 0.42, [[(label, 12, True, fg, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide(); bg(s, NAVY)
box(s, 0, 0, 13.333, 0.14, fill=TEAL)
box(s, 0, 7.36, 13.333, 0.14, fill=TEAL)
text(s, 1.0, 2.15, 11.3, 2.2, [
    [("OWASP Top 10", 54, True, WHITE, FONT)],
    [("2025", 54, True, TEAL, FONT)],
], space_after=2)
text(s, 1.05, 4.05, 11.3, 1.0, [
    [("The ten most critical web application security risks — ", 22, False, LIGHT, FONT),
     ("explained, broken, and fixed in an hour.", 22, True, WHITE, FONT)],
])
text(s, 1.05, 5.5, 11.3, 1.0, [
    [("A hands-on session for engineers.  Slides + a companion repo you run on your laptop.", 16, False, RGBColor(0x9F,0xB0,0xC8), FONT)],
    [("Source: owasp.org/Top10/2025", 14, False, TEAL, FONT)],
], space_after=6)

# ===========================================================================
# 2. WHY THIS HOUR
# ===========================================================================
s = slide(); bg(s, WHITE)
header(s, "", "Why this hour", accent=TEAL)
text(s, 0.55, 1.15, 12.2, 0.9, [
    [("Security isn't a separate job — it's a property of the code you already write. ",
      19, False, INK, FONT),
     ("This hour makes the ten most common failure modes concrete.", 19, True, INK, FONT)],
], anchor=MSO_ANCHOR.MIDDLE)

items = [
    ("See it", "Every category has a working exploit you can run in one command.", RED),
    ("Understand it", "The broken code and the fix sit side by side, with comments.", AMBER),
    ("Fix it", "The fix is shown in the smallest honest form — the idea, not a library.", GREEN),
    ("Prove it", "A test exploits the flaw, then proves the fix holds. 74 tests, all green.", TEAL),
]
x = 0.55; w = 2.95; gap = 0.18
for i,(t,d,c) in enumerate(items):
    xi = x + i*(w+gap)
    box(s, xi, 2.4, w, 3.2, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    box(s, xi, 2.4, w, 0.7, fill=c, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, xi, 2.4, w, 0.7, [[(t, 18, True, WHITE, FONT)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, xi+0.22, 3.3, w-0.44, 2.1, [[(d, 15, False, INK, FONT)]], line_spacing=1.1)
text(s, 0.55, 6.1, 12.2, 0.8, [
    [("The mental model for all ten:  ", 17, True, INK, FONT),
     ("untrusted input crossing a trust boundary.", 17, False, GREY, FONT)],
], anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# 3. HOW THE LAB WORKS
# ===========================================================================
s = slide(); bg(s, WHITE)
header(s, "", "How the companion repo works", accent=TEAL)
text(s, 0.55, 1.35, 12.2, 0.7, [
    [("Every category is mounted twice. Send the same request to both and diff the response.", 18, False, INK, FONT)],
])
chip(s, 0.55, 2.15, 3.0, "/vuln/aNN   — broken", RED)
chip(s, 3.85, 2.15, 3.0, "/safe/aNN   — fixed", GREEN)
codebox(s, 0.55, 2.9, 12.2, 2.15, [
    ("$ npm install && npm run doctor && npm start", CODEFG),
    ("  OWASP Top 10:2025 lab -> http://localhost:3000", RGBColor(0x8A,0x99,0xAD)),
    ("", CODEFG),
    ("$ npm test", CODEFG),
    ("  # tests 74   # pass 74   # fail 0", TEAL),
], title="terminal")
text(s, 0.55, 5.3, 12.2, 1.6, [
    [("For each category:  ", 16, True, INK, FONT),
     ("labs/aNN/README.md", 16, False, GREEN, MONO),
     ("  walks you through ", 16, False, INK, FONT),
     ("Break it -> Read it -> Fix it -> Prove it,", 16, True, INK, FONT),
     ("  and ends with a ", 16, False, INK, FONT),
     ("“spot it in code review”", 16, True, INK, FONT),
     (" checklist you can use at work on Monday.", 16, False, INK, FONT)],
    [("Deliberately vulnerable — it binds to localhost only. Never deploy it or point it at real data.",
      14, False, RED, FONT)],
], space_after=10)

# ===========================================================================
# 4. HOW THE LIST IS BUILT
# ===========================================================================
s = slide(); bg(s, WHITE)
header(s, "", "How the list is built", accent=TEAL, kicker="Data-informed, not blindly data-driven")
stats = [
    ("2.8M", "applications analysed"),
    ("~175k", "CVE records (up from 125k)"),
    ("643", "unique CWEs mapped"),
    ("8 + 2", "from data + community survey"),
]
x=0.55; w=2.95; gap=0.18
for i,(n,d) in enumerate(stats):
    xi = x+i*(w+gap)
    box(s, xi, 1.55, w, 1.7, fill=NAVY, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, xi, 1.7, w, 0.9, [[(n, 34, True, TEAL, FONT)]], align=PP_ALIGN.CENTER)
    text(s, xi+0.15, 2.55, w-0.3, 0.6, [[(d, 13, False, LIGHT, FONT)]], align=PP_ALIGN.CENTER, line_spacing=1.0)
text(s, 0.55, 3.7, 12.2, 3.4, [
    [("Eight categories come from the contributed data — how often each weakness actually appears in tested apps.", 17, False, INK, FONT)],
    [("Two come from a community survey of AppSec practitioners, to capture risks the data can't see yet.", 17, False, INK, FONT)],
    [("", 8, False, INK, FONT)],
    [("Why the survey?  ", 17, True, INK, FONT),
     ("Testing data looks backwards. It takes years for a new class of attack to become something tools "
      "can detect at scale. The survey asks the people on the front line what they're seeing now.", 17, False, GREY, FONT)],
    [("", 8, False, INK, FONT)],
    [("Survey-driven this year:  ", 16, True, INK, FONT),
     ("A03 Software Supply Chain Failures", 16, True, AMBER, FONT),
     ("  and  ", 16, False, GREY, FONT),
     ("A09 Security Logging & Alerting Failures.", 16, True, AMBER, FONT)],
], space_after=8, line_spacing=1.08)

# ===========================================================================
# 5. WHAT CHANGED 2021 -> 2025
# ===========================================================================
s = slide(); bg(s, WHITE)
header(s, "", "What changed from 2021", accent=TEAL)
rows = [
    ("A01", "Broken Access Control", "#1  →  #1", "unchanged; SSRF folded in", GREY),
    ("A02", "Security Misconfiguration", "#5  →  #2", "biggest mover", GREEN),
    ("A03", "Software Supply Chain Failures", "NEW", "expanded from “Vulnerable Components”; #1 in survey", AMBER),
    ("A04", "Cryptographic Failures", "#2  →  #4", "down two", RED),
    ("A05", "Injection", "#3  →  #5", "down two; XSS lives here", RED),
    ("A06", "Insecure Design", "#4  →  #6", "down two", RED),
    ("A07", "Authentication Failures", "#7  →  #7", "renamed", GREY),
    ("A08", "Software or Data Integrity Failures", "#8  →  #8", "unchanged", GREY),
    ("A09", "Security Logging & Alerting Failures", "#9  →  #9", "“Monitoring” → “Alerting”", GREY),
    ("A10", "Mishandling of Exceptional Conditions", "NEW", "brand new category", AMBER),
]
y = 1.45; rh = 0.56
for i,(code,name,delta,note,c) in enumerate(rows):
    yi = y + i*rh
    if i % 2 == 0:
        box(s, 0.55, yi, 12.23, rh, fill=LIGHT)
    text(s, 0.7, yi, 1.0, rh, [[(code, 15, True, NAVY, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.7, yi, 5.4, rh, [[(name, 15, False, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 7.2, yi, 1.9, rh, [[(delta, 14, True, c, MONO)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 9.2, yi, 3.5, rh, [[(note, 13, False, GREY, FONT)]], anchor=MSO_ANCHOR.MIDDLE)

# ===========================================================================
# 6. AGENDA
# ===========================================================================
s = slide(); bg(s, NAVY)
box(s, 0, 1.15, 13.333, 0.06, fill=TEAL)
text(s, 0.55, 0.25, 12, 0.9, [[("The next hour", 30, True, WHITE, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
plan = [
    ("0:00", "Framing — how to think about all ten", False),
    ("0:05", "A01  Broken Access Control", True),
    ("0:13", "A05  Injection", True),
    ("0:21", "A02  Security Misconfiguration", True),
    ("0:29", "A03  Software Supply Chain Failures", True),
    ("0:37", "Rapid survey — A04, A06, A07, A08, A09, A10", False),
    ("0:55", "What to do on Monday", False),
]
y=1.55; rh=0.75
for t,label,lab in plan:
    box(s, 0.9, y+0.06, 1.1, 0.5, fill=TEAL if lab else NAVY2, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 0.9, y+0.06, 1.1, 0.5, [[(t, 15, True, NAVY if lab else LIGHT, MONO)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.2, y, 9.0, 0.62, [[(label, 19, lab, WHITE if lab else LIGHT, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    if lab:
        text(s, 11.2, y, 1.6, 0.62, [[("LAB", 13, True, TEAL, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    y += rh

# ===========================================================================
# Category slide builders
# ===========================================================================
def cat_break(code, title, rank, oneliner, break_lines, scenarios, accent=TEAL, code_title="terminal"):
    s = slide(); bg(s, WHITE)
    header(s, code, title, accent=accent, kicker=rank)
    text(s, 0.55, 1.35, 12.2, 0.85, [
        [(oneliner, 18, True, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    # left: break it
    chip(s, 0.55, 2.35, 2.2, "BREAK IT", RED)
    codebox(s, 0.55, 2.95, 7.0, 3.9, break_lines, title=code_title)
    # right: what's really happening
    chip(s, 7.9, 2.35, 3.0, "WHAT'S HAPPENING", GREY)
    runs = [[(sc, 15, False, INK, FONT)] for sc in scenarios]
    box(s, 7.9, 2.95, 4.88, 3.9, fill=REDBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    text(s, 8.15, 3.15, 4.4, 3.5, runs, space_after=10, line_spacing=1.08)
    return s

def cat_fix(code, title, fix_points, spot_points, accent=GREEN):
    s = slide(); bg(s, WHITE)
    header(s, code, title + " — the fix", accent=accent)
    chip(s, 0.55, 1.45, 2.2, "FIX IT", GREEN)
    box(s, 0.55, 2.05, 7.0, 4.75, fill=GRNBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs = []
    for p in fix_points:
        runs.append([("•  ", 16, True, GREEN, FONT), (p[0], 16, True, INK, FONT)])
        if len(p) > 1:
            runs.append([("    "+p[1], 14, False, GREY, FONT)])
    text(s, 0.8, 2.25, 6.5, 4.4, runs, space_after=7, line_spacing=1.06)
    chip(s, 7.9, 1.45, 3.4, "SPOT IT IN REVIEW", NAVY)
    box(s, 7.9, 2.05, 4.88, 4.75, fill=LIGHT, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs2 = [[("⚑  ", 14, True, RED, FONT), (p, 14.5, False, INK, FONT)] for p in spot_points]
    text(s, 8.15, 2.3, 4.4, 4.3, runs2, space_after=9, line_spacing=1.08)
    return s

def cat_survey(code, title, rank, oneliner, break_line, fix_points, accent=TEAL):
    s = slide(); bg(s, WHITE)
    header(s, code, title, accent=accent, kicker=rank)
    text(s, 0.55, 1.35, 12.2, 0.8, [[(oneliner, 17, True, INK, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
    chip(s, 0.55, 2.3, 2.2, "BREAK IT", RED)
    codebox(s, 0.55, 2.9, 6.0, 3.9, break_line, title="try it")
    chip(s, 6.95, 2.3, 2.2, "FIX IT", GREEN)
    box(s, 6.95, 2.9, 5.83, 3.9, fill=GRNBG, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    runs = [[("•  ", 15, True, GREEN, FONT), (p, 15, False, INK, FONT)] for p in fix_points]
    text(s, 7.2, 3.12, 5.4, 3.5, runs, space_after=9, line_spacing=1.1)
    return s

C = CODEFG; G = RGBColor(0x8A,0x99,0xAD); TL = TEAL; RD = RGBColor(0xF2,0x8B,0x8B)

# ---- A01 -------------------------------------------------------------------
cat_break("A01:2025", "Broken Access Control", "#1 — unchanged. Found in nearly every app tested.",
    "Enforcing that users can only do what they're allowed to.",
    [("# Alice reads Carol's confidential invoice —", G),
     ("# just by changing the id in the URL.", G),
     ("$ curl .../vuln/a01/invoices/1004?as=alice", C),
     ('{"invoice":{"user_id":3,', RD),
     ('  "memo":"SEVERANCE PAYMENT..."}}', RD),
     ("", C),
     ("# The fixed route scopes the query to Alice:", G),
     ("$ curl .../safe/a01/invoices/1004?as=alice", C),
     ("HTTP 404", TL)],
    ["IDOR: the record id comes from the request and is trusted — no ownership check.",
     "Forced browsing: the admin page's only “protection” is that it isn't linked.",
     "Client-side control: the UI hides the button; curl doesn't care."],
    accent=TEAL)
cat_fix("A01:2025", "Broken Access Control",
    [("Deny by default", "for anything non-public."),
     ("Enforce ownership in the query", "…AND user_id = ? — no path can forget it."),
     ("Check the role on the server", "regardless of what the UI shows."),
     ("Reuse one mechanism", "not a hand-rolled check per endpoint."),
     ("Return 404, not 403", "a 403 confirms the id exists.")],
    ["An id/filename from the request used to fetch a record with no ownership filter",
     "Authorization logic that lives only in front-end code",
     "if (user.isAdmin) in the template but not the handler",
     "SSRF (now in A01): server fetches a user-supplied URL"])

# ---- A05 -------------------------------------------------------------------
cat_break("A05:2025", "Injection", "#5 (was #3). Most CVEs of any category. XSS lives here.",
    "Untrusted input reaches an interpreter as CODE instead of DATA.",
    [("# Log in as admin with no password:", G),
     ("$ curl .../vuln/a05/login \\", C),
     ("   -d \"username=admin'--\" -d pass=x", C),
     ('{"loggedIn":true,"user":"admin"}', RD),
     ("", C),
     ("# Same payload, parameterised query:", G),
     ("$ curl .../safe/a05/login ...", C),
     ('{"loggedIn":false}', TL),
     ("", C),
     ("# Command injection: host=localhost;id", RD)],
    ["SQL injection: the payload becomes part of the query's syntax.",
     "UNION attacks read arbitrary tables (e.g. the password hashes).",
     "OS command injection: a ; runs a second shell command.",
     "XSS: input echoed into HTML runs as script in the victim's browser."],
    accent=TEAL)
cat_fix("A05:2025", "Injection",
    [("Parameterised queries", "? placeholders — the #1 habit in this workshop."),
     ("No shell", "execFile(cmd, [args]) not exec(string)."),
     ("Positive (allow-list) validation", "of shape, on the server."),
     ("Encode output for its context", "escapeHtml for HTML — the fix for XSS."),
     ("Content-Security-Policy", "as defence in depth.")],
    ["SQL built with + , template literals, or .format()",
     "exec(), system(), eval() with anything from a request",
     "User input written into HTML without encoding",
     "“We validate on the client” as the only defence"])

# ---- A02 -------------------------------------------------------------------
cat_break("A02:2025", "Security Misconfiguration", "▲ #5 → #2. 100% of tested apps had some form.",
    "The code may be fine; the SETTINGS around it hand attackers a start.",
    [("# Directory listing exposes a forgotten file", G),
     ("$ curl .../vuln/a02/files/", C),
     (" .env.backup", RD),
     ("$ curl .../vuln/a02/files/.env.backup", C),
     (" STRIPE_KEY=sk_live_51Hxxxx", RD),
     (" DATABASE_URL=postgres://app:s3cr3t@...", RD),
     ("", C),
     ("# An error returns a full stack trace:", G),
     ("$ curl .../vuln/a02/boom", C),
     (" Error: db.internal:5432 refused ...", RD)],
    ["Directory listing + a file nobody meant to ship = leaked secrets.",
     "Verbose errors reveal internal hosts and component versions.",
     "Default accounts and debug endpoints left enabled.",
     "Cookies with no HttpOnly / Secure / SameSite."],
    accent=GREEN)
cat_fix("A02:2025", "Security Misconfiguration",
    [("Repeatable hardening", "so every env ships locked down."),
     ("Serve only what you intend", "no listing, no dotfiles."),
     ("No debug endpoints", "in the deployed artifact."),
     ("Opaque errors out", "detail to the log with a reference id."),
     ("Cookie attributes + security headers", "HttpOnly/Secure/SameSite; helmet.")],
    ["dotfiles:'allow' or directory listing on a static mount",
     "res.send(err.stack) or returning err.message",
     "A /debug or /actuator endpoint with no auth",
     "Cookies set with no options"])

# ---- A03 -------------------------------------------------------------------
cat_break("A03:2025", "Software Supply Chain Failures", "NEW scope · voted #1 concern in the survey.",
    "You ship a lot of code you didn't write — and the pipeline that builds it.",
    [("# We declare 2 dependencies.", G),
     ("# How many actually get installed?", G),
     ("$ curl .../vuln/a03/inventory", C),
     ('{"directDependencies":2,', C),
     ('  "packagesActuallyInstalled":  57,', RD),
     ('  "transitiveMultiplier": 28.5}', RD),
     ("", C),
     ("$ npm audit          # known CVEs", C),
     ("$ npm run sbom       # bill of materials", C),
     ("$ npm ci  (not install)  in CI", TL)],
    ["Expanded from 2021's “Vulnerable & Outdated Components”.",
     "Now covers dependencies AND their transitive deps, build systems, CI/CD, and tooling.",
     "SolarWinds; Log4Shell; the 2025 Shai-Hulud npm worm (self-propagating via install scripts).",
     "Fewest data occurrences, highest exploit + impact scores."],
    accent=AMBER)
cat_fix("A03:2025", "Software Supply Chain Failures",
    [("Generate & keep an SBOM", "your “are we affected?” answer sheet."),
     ("Track transitive deps", "not just the ones you named."),
     ("Commit the lockfile, use npm ci", "installs exact, integrity-checked versions."),
     ("--ignore-scripts by default", "install scripts spread the 2025 worm."),
     ("Pin Actions to a commit SHA", "not a mutable @v4 tag."),
     ("Separation of duties", "no one writes + promotes unreviewed.")],
    ["npm install in CI or a Dockerfile (should be npm ci)",
     "A missing or gitignored lockfile",
     "GitHub Actions pinned to floating tags",
     "Security-critical dep on a caret range, no lockfile"],
    accent=AMBER)

# ---- Survey divider --------------------------------------------------------
s = slide(); bg(s, NAVY)
box(s, 0, 3.3, 13.333, 0.06, fill=TEAL)
text(s, 0.55, 2.3, 12.2, 1.0, [[("Rapid survey", 40, True, WHITE, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.55, 3.5, 12.2, 1.2, [
    [("A04 · A06 · A07 · A08 · A09 · A10", 22, True, TEAL, FONT)],
    [("One demo and one fix each. All six are full labs in the repo.", 16, False, LIGHT, FONT)],
], space_after=10)

# ---- A04 survey ------------------------------------------------------------
cat_survey("A04:2025", "Cryptographic Failures", "▼ #2 → #4. Weak randomness dominates.",
    "Crypto that's absent, weak, unsalted, predictable, or unauthenticated.",
    [("# reset tokens from Math.random():", G),
     ("$ curl .../vuln/a04/reset-token", C),
     ('{"token":"a1b2c3d4"}', RD),
     ("# unsalted MD5 — look it up, no", G),
     ("# cracking needed:", G),
     ("21232f29...  = md5(\"admin\")", RD),
     ("", C),
     ("# AES-CBC w/o a MAC: attacker", G),
     ("# rewrites \"user\" -> \"root\",", G),
     ("# no key required.", RD)],
    ["crypto.randomBytes (a CSPRNG), never Math.random()",
     "Store passwords with a salted, slow KDF (scrypt / argon2 / bcrypt)",
     "Authenticated encryption (AES-GCM): tampering fails hard",
     "TLS 1.2+ only; drop MD5 / SHA-1"],
    accent=RED)

# ---- A06 survey ------------------------------------------------------------
cat_survey("A06:2025", "Insecure Design", "▼ #4 → #6. A flaw in the plan, not the code.",
    "Missing control DESIGN — you can't patch it in later.",
    [("# 'One per order' coupon, stacked", G),
     ("# ten times — each use is valid:", G),
     ("$ curl .../vuln/a06/checkout \\", C),
     ('  -d coupons=[x10]', C),
     ('{"totalCents": 3487}', RD),
     ("", C),
     ("# refund more than was ever paid;", G),
     ("# 'mother's maiden name' = Smith", G),
     ("#   -> reset token.", RD)],
    ["State the rule, then enforce it server-side",
     "Put invariants next to the operation",
     "Threat-model the critical flows before building",
     "Recovery via a channel the user controls — no security questions"],
    accent=RED)

# ---- A07 survey ------------------------------------------------------------
cat_survey("A07:2025", "Authentication Failures", "#7 — renamed. It's the stuff AROUND the password.",
    "Weak throttling, no MFA, default creds, broken sessions, reused breaches.",
    [("$ node scripts/\\", C),
     ("  credential-stuffing.js admin", C),
     (" /vuln: got 'admin'  (default!)", RD),
     (" /safe: locked out + blocked", TL),
     ("", C),
     ("# session fixation: the id the", G),
     ("# attacker planted survives login", G),
     ("# on /vuln — not on /safe.", G)],
    ["MFA — OWASP's first prevention bullet here",
     "No default credentials; block known-breached passwords",
     "Rate-limit / lock out (per-account and per-IP)",
     "Rotate the session id at every privilege change"],
    accent=GREY)

# ---- A08 survey ------------------------------------------------------------
cat_survey("A08:2025", "Software or Data Integrity Failures", "#8 — the low-level cousin of A03.",
    "Did this artifact / data really come from who I think, unaltered?",
    [("# Prototype pollution: one request", G),
     ("# makes EVERY object isAdmin.", G),
     ("$ curl .../vuln/a08/preferences \\", C),
     ('  -d \'{"__proto__":{"isAdmin":1}}\'', C),
     ('{"pollutedBystander":', RD),
     ('  {"isAdmin":true}}', RD),
     ("", C),
     ("# unsigned 'auto-update' applied", G),
     ("# from any source.", G)],
    ["Reject __proto__ / constructor; allow-list known keys",
     "Verify a signature from the expected source first",
     "Don't deserialize untrusted data into live objects",
     "Review process + CI/CD segregation"],
    accent=GREY)

# ---- A09 survey ------------------------------------------------------------
cat_survey("A09:2025", "Security Logging & Alerting Failures", "#9 — “Monitoring” → “Alerting”.",
    "Great logging with no ALERTING can't find an incident.",
    [("# 50 failed logins on /vuln...", G),
     ("$ curl .../vuln/a09/logs", C),
     (" <-- empty. No trace at all.", RD),
     ("", C),
     ("# log injection: a newline in a", G),
     ("# username forges a whole entry.", G),
     ("# whole request bodies (with", G),
     ("# passwords) dumped to the log.", RD)],
    ["Log every auth / access-control / validation failure",
     "Structured (JSON) logging neutralises log injection",
     "Redact secrets; never dump request bodies",
     "Alert on thresholds — the half the rename points at"],
    accent=GREY)

# ---- A10 survey ------------------------------------------------------------
cat_survey("A10:2025", "Mishandling of Exceptional Conditions", "BRAND NEW. What happens when things go wrong.",
    "Failing OPEN, leaking internals, partial state when step 3 of 4 throws.",
    [("# authz service is down ->", G),
     ("# the catch block says 'yes':", G),
     ("$ curl -X POST \\", C),
     ("  .../vuln/a10/admin-action", C),
     ('{"ranAdminAction":true}  # !!', RD),
     ("", C),
     ("# non-atomic transfer: debit", G),
     ("# succeeds, credit throws,", G),
     ("# money vanishes.", RD)],
    ["Fail CLOSED — if 'may they?' is unknown, deny (503)",
     "Catch errors where they happen; global handler as backstop",
     "Opaque errors out, detail to the log",
     "One transaction — roll back every step on failure"],
    accent=AMBER)

# ===========================================================================
# WRAP — MONDAY
# ===========================================================================
s = slide(); bg(s, WHITE)
header(s, "", "What to do on Monday", accent=TEAL)
text(s, 0.55, 1.3, 12.2, 0.7, [[("Pick one. Small and shipped beats big and someday.", 19, True, INK, FONT)]],
     anchor=MSO_ANCHOR.MIDDLE)
todos = [
    ("Turn on a linter rule", "eslint-plugin-security or your language's equivalent catches injection and unsafe APIs for free.", TEAL),
    ("Commit a lockfile + npm ci", "Pin your dependencies and switch CI off npm install. One PR, permanent A03 win.", AMBER),
    ("Add a breached-password check", "Have-I-Been-Pwned's k-anonymity API, on registration and password change.", GREEN),
    ("Use the review checklists", "Each labs/aNN/README.md ends with “spot it in review”. Bring them to your next PR.", RED),
]
y=2.25; h=1.05; gap=0.16
for t,d,c in todos:
    box(s, 0.55, y, 0.14, h, fill=c)
    box(s, 0.69, y, 12.09, h, fill=LIGHT)
    text(s, 0.95, y, 11.6, h, [
        [(t, 18, True, INK, FONT)],
        [(d, 14.5, False, GREY, FONT)],
    ], anchor=MSO_ANCHOR.MIDDLE, space_after=2, line_spacing=1.05)
    y += h+gap

# ===========================================================================
# RESOURCES / CLOSE
# ===========================================================================
s = slide(); bg(s, NAVY)
box(s, 0, 0, 13.333, 0.14, fill=TEAL)
text(s, 0.9, 1.4, 11.5, 1.0, [[("Take it further", 34, True, WHITE, FONT)]])
text(s, 0.95, 2.7, 11.5, 3.6, [
    [("The repo", 20, True, TEAL, FONT)],
    [("Ten full labs, 74 tests, exploit scripts, and a facilitator guide. Break every category, "
      "not just the four we covered.", 16, False, LIGHT, FONT)],
    [("", 8, False, LIGHT, FONT)],
    [("The source", 20, True, TEAL, FONT)],
    [("owasp.org/Top10/2025  — every category page has real stats, attack scenarios, "
      "and a full “How to Prevent” list.", 16, False, LIGHT, FONT)],
    [("", 8, False, LIGHT, FONT)],
    [("Beyond the Top 10", 20, True, TEAL, FONT)],
    [("OWASP ASVS (verification standard) · the OWASP Cheat Sheet Series · your language's "
      "secure-coding guide.", 16, False, LIGHT, FONT)],
], space_after=8, line_spacing=1.08)
text(s, 0.95, 6.6, 11.5, 0.6, [[("Security is a habit, not a phase. Thanks for the hour.", 16, True, WHITE, FONT)]])

import os
out = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "docs", "OWASP_Top_10_2025.pptx")
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
